from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional
import os
from dotenv import load_dotenv
import os
from openai import OpenAI
import base64
from PIL import Image
import io
import PyPDF2
from docx import Document
from pptx import Presentation
import json

# Load environment variables
load_dotenv()

# Initialize FastAPI app
app = FastAPI(title="CUSTOS AI Service", version="1.0.0")

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],  # Next.js dev server
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure OpenAI
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
if OPENAI_API_KEY:
    client = OpenAI(api_key=OPENAI_API_KEY)
else:
    client = None

# Response models
class Student(BaseModel):
    name: str
    roll_number: Optional[str] = None

class ExtractRosterResponse(BaseModel):
    students: List[Student]
    total_count: int
    success: bool
    message: str

# Health check endpoint
@app.get("/")
async def root():
    return {
        "service": "CUSTOS AI Service",
        "version": "1.0.0",
        "status": "running",
        "ai_provider": "OpenAI GPT-4o-mini",
        "openai_configured": OPENAI_API_KEY is not None
    }

@app.get("/health")
async def health_check():
    return {
        "status": "healthy",
        "openai_api": "configured" if OPENAI_API_KEY else "not configured"
    }

# Vision AI endpoint for student roster extraction
@app.post("/api/vision/extract-roster", response_model=ExtractRosterResponse)
async def extract_roster(file: UploadFile = File(...)):
    """
    Extract student names from an uploaded image using OpenAI GPT-4o-mini Vision.
    
    Accepts: JPG, PNG, JPEG images
    Returns: List of student names extracted from the image
    """
    
    if not client:
        raise HTTPException(
            status_code=500,
            detail="OpenAI API is not configured. Please set OPENAI_API_KEY in .env file"
        )
    
    # Validate file type
    if not file.content_type or not file.content_type.startswith("image/"):
        raise HTTPException(
            status_code=400,
            detail="File must be an image (JPG, PNG, JPEG)"
        )
    
    try:
        # Read and encode the image
        contents = await file.read()
        base64_image = base64.b64encode(contents).decode('utf-8')
        
        # Prepare the prompt for OpenAI
        prompt = """Analyze this image which contains a list of student names. 
Extract all student names from the image.

Rules:
- Extract ONLY the names, ignore any other text
- If there are roll numbers or IDs, include them
- Handle handwritten or printed text
- Ignore headers, dates, and other metadata

Return ONLY a JSON object in this exact format:
{
    "students": [
        {"name": "Student Name 1", "roll_number": "1"},
        {"name": "Student Name 2", "roll_number": "2"}
    ]
}

If no names are found, return: {"students": []}"""
        
        # Call OpenAI Vision API
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=1000
        )
        
        response_text = response.choices[0].message.content
        
        # Parse the JSON response
        import json
        import re
        
        # Extract JSON from markdown code blocks if present
        json_match = re.search(r'```json\s*(.*?)\s*```', response_text, re.DOTALL)
        if json_match:
            json_str = json_match.group(1)
        else:
            # Try to find JSON object in the response
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            if json_match:
                json_str = json_match.group(0)
            else:
                json_str = response_text
        
        # Parse JSON
        try:
            data = json.loads(json_str)
            students = data.get("students", [])
            
            # Convert to Student models
            student_list = [Student(**student) for student in students]
            
            return ExtractRosterResponse(
                students=student_list,
                total_count=len(student_list),
                success=True,
                message=f"Successfully extracted {len(student_list)} student(s)"
            )
        except json.JSONDecodeError:
            # Fallback: Try to extract names as plain text lines
            lines = response_text.strip().split('\n')
            students = []
            for line in lines:
                line = line.strip()
                if line and len(line) > 2:  # Filter out empty or very short lines
                    # Remove numbering (1., 2., etc.)
                    name = re.sub(r'^\d+[\.\)]\s*', '', line)
                    if name:
                        students.append(Student(name=name))
            
            return ExtractRosterResponse(
                students=students,
                total_count=len(students),
                success=True,
                message=f"Extracted {len(students)} student(s) (fallback mode)"
            )
            
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Error processing image: {str(e)}"
        )

# Syllabus extraction endpoint
@app.post("/api/syllabus/extract")
async def extract_syllabus(file: UploadFile = File(...)):
    """
    Extract content from syllabus documents (PDF, DOCX, PPTX, Images)
    Returns structured JSON with 99% storage optimization
    """
    if not client:
        raise HTTPException(status_code=500, detail="OpenAI API not configured")
    
    try:
        # Read file content
        content = await file.read()
        file_type = file.content_type
        
        # Extract text based on file type
        extracted_text = ""
        
        if "pdf" in file_type:
            # Extract from PDF
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            for page in pdf_reader.pages:
                extracted_text += page.extract_text() + "\n\n"
                
        elif "wordprocessingml" in file_type:  # DOCX
            # Extract from Word document
            doc_file = io.BytesIO(content)
            doc = Document(doc_file)
            for paragraph in doc.paragraphs:
                extracted_text += paragraph.text + "\n"
                
        elif "presentationml" in file_type:  # PPTX
            # Extract from PowerPoint
            ppt_file = io.BytesIO(content)
            prs = Presentation(ppt_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
                        
        elif "image" in file_type:
            # Use OpenAI Vision for images
            base64_image = base64.b64encode(content).decode('utf-8')
            
            vision_response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Extract all text content from this syllabus/textbook page. Include headings, paragraphs, formulas, examples, and any other text visible."
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{file_type};base64,{base64_image}"
                            }
                        }
                    ]
                }],
                max_tokens=2000
            )
            
            extracted_text = vision_response.choices[0].message.content
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")
        
        # Use AI to structure the content
        structure_prompt = f"""
        Analyze this syllabus/textbook content and extract structured information.
        
        Content:
        {extracted_text[:8000]}  # Limit to avoid token limits
        
        Return a JSON object with the following structure:
        {{
            "title": "Chapter title or main topic",
            "sections": [
                {{
                    "heading": "Section heading",
                    "text": "Section content",
                    "page": 1
                }}
            ],
            "formulas": ["formula1", "formula2"],
            "key_points": ["point1", "point2"],
            "examples": [
                {{
                    "question": "Example question",
                    "solution": "Solution steps",
                    "answer": "Final answer"
                }}
            ],
            "definitions": ["term: definition"]
        }}
        
        Be comprehensive but concise. Extract all important information.
        """
        
        structure_response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{
                "role": "user",
                "content": structure_prompt
            }],
            temperature=0.3,
            max_tokens=3000
        )
        
        # Parse structured content
        try:
            structured_content = json.loads(structure_response.choices[0].message.content)
        except json.JSONDecodeError:
            # Fallback if AI doesn't return valid JSON
            structured_content = {
                "title": "Extracted Content",
                "sections": [{
                    "heading": "Content",
                    "text": extracted_text[:1000],
                    "page": 1
                }],
                "formulas": [],
                "key_points": [],
                "examples": [],
                "definitions": []
            }
        
        return structured_content
        
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Error extracting content: {str(e)}"
        )

# ==========================================
# LESSON PLAN GENERATION MODELS
# ==========================================

class Topic(BaseModel):
    topic_id: str
    topic_title: str
    duration_minutes: int
    difficulty: str
    learning_objectives: Optional[List[str]] = []

class Constraints(BaseModel):
    total_days: int
    periods_per_week: int
    period_duration_minutes: int
    holidays: List[str] = []

class LessonPlanRequest(BaseModel):
    topics: List[Topic]
    constraints: Constraints

# ==========================================
# LESSON PLAN GENERATION ENDPOINT
# ==========================================

@app.post("/api/lesson-plan/generate")
async def generate_lesson_plan(request: LessonPlanRequest):
    """
    Generate an optimal day-by-day lesson plan based on topics and constraints.
    """
    if not client:
        raise HTTPException(status_code=500, detail="OpenAI API not configured")
    
    try:
        # Construct the prompt
        topics_str = json.dumps([t.dict() for t in request.topics], indent=2)
        constraints_str = json.dumps(request.constraints.dict(), indent=2)
        
        prompt = f"""
        You are an expert curriculum planner. Create an optimal day-by-day lesson plan.
        
        TOPICS TO COVER:
        {topics_str}
        
        CONSTRAINTS:
        {constraints_str}
        
        REQUIREMENTS:
        1. Distribute topics evenly across the available days ({request.constraints.total_days} days).
        2. Consider difficulty progression (easy → hard).
        3. Include review interactions/activities.
        4. Balance theory and practice.
        5. Respect the total duration constraints.
        
        Return ONLY a JSON object with this EXACT structure:
        {{
            "schedule": [
                {{
                    "day": 1,
                    "date": "YYYY-MM-DD",  # Use relative dates from start (or placeholders if no start date)
                    "topic_id": "uuid_from_input",
                    "topic_title": "String",
                    "activities": ["Activity 1", "Activity 2"],
                    "duration": 45
                }}
            ],
            "summary": {{
                "total_days": 30,
                "topics_covered": 10,
                "percent_utilization": 85
            }}
        }}
        """
        
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{
                "role": "user",
                "content": prompt
            }],
            temperature=0.3, # Low temperature for structured output
            response_format={"type": "json_object"}
        )
        
        return json.loads(response.choices[0].message.content)
        
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Error generating lesson plan: {str(e)}"
        )


# ==========================================
# AI RESOURCE GENERATION
# ==========================================

class ResourceGenerationRequest(BaseModel):
    topic_id: str
    topic_title: str
    topic_content: dict  # Content from lesson_topics.content
    resource_type: str  # lesson_notes, study_guide, worksheet, revision_notes, formulas_list
    grade_level: int
    subject_name: str

@app.post("/api/resources/generate")
async def generate_resource(request: ResourceGenerationRequest):
    """
    Generate AI-powered educational resources for a topic.
    Supports 5 resource types: lesson_notes, study_guide, worksheet, revision_notes, formulas_list
    """
    if not client:
        raise HTTPException(status_code=500, detail="OpenAI API not configured")
    
    try:
        # Build prompt based on resource type
        type_instructions = {
            "lesson_notes": """
                Create detailed LESSON NOTES for teaching this topic:
                - Clear explanations suitable for the grade level
                - Step-by-step breakdowns of concepts
                - Examples with solutions
                - Teaching tips and common misconceptions
                - Estimated time for each section
                
                Format as JSON:
                {
                    "title": "Topic Title",
                    "sections": [
                        {"heading": "...", "content": "...", "duration_minutes": 10}
                    ],
                    "examples": [{"question": "...", "solution": "..."}],
                    "teaching_tips": ["..."],
                    "common_mistakes": ["..."]
                }
            """,
            "study_guide": """
                Create a STUDY GUIDE for students to review this topic:
                - Key concepts summary
                - Important definitions
                - Quick reference points
                - Memory aids and mnemonics
                - Self-assessment questions
                
                Format as JSON:
                {
                    "title": "Topic Title - Study Guide",
                    "key_concepts": ["..."],
                    "definitions": [{"term": "...", "meaning": "..."}],
                    "summary": "...",
                    "memory_tips": ["..."],
                    "quick_quiz": [{"question": "...", "answer": "..."}]
                }
            """,
            "worksheet": """
                Create a PRACTICE WORKSHEET for this topic:
                - 10-15 practice problems of varying difficulty
                - Include multiple choice, short answer, and problem-solving
                - Provide answer key separately
                - Progressive difficulty (easy → medium → hard)
                
                Format as JSON:
                {
                    "title": "Topic Title - Practice Worksheet",
                    "instructions": "...",
                    "problems": [
                        {"number": 1, "type": "multiple_choice", "question": "...", "options": ["A", "B", "C", "D"], "difficulty": "easy"}
                    ],
                    "answer_key": [{"number": 1, "answer": "...", "explanation": "..."}]
                }
            """,
            "revision_notes": """
                Create ULTRA-CONDENSED REVISION NOTES (cheat sheet) for this topic:
                - Maximum 1 page of content
                - Only the most critical formulas and facts
                - Bullet points only
                - Perfect for last-minute review
                
                Format as JSON:
                {
                    "title": "Topic Title - Quick Revision",
                    "must_remember": ["..."],
                    "key_formulas": ["..."],
                    "important_facts": ["..."],
                    "exam_tips": ["..."]
                }
            """,
            "formulas_list": """
                Create a FORMULAS & DEFINITIONS reference for this topic:
                - All formulas with variable explanations
                - All definitions with examples
                - Units and conversions if applicable
                - When to use each formula
                
                Format as JSON:
                {
                    "title": "Topic Title - Formulas & Definitions",
                    "formulas": [
                        {"formula": "...", "name": "...", "variables": {"x": "explanation"}, "when_to_use": "..."}
                    ],
                    "definitions": [{"term": "...", "definition": "...", "example": "..."}],
                    "units": [{"quantity": "...", "unit": "...", "symbol": "..."}]
                }
            """
        }
        
        instruction = type_instructions.get(request.resource_type, type_instructions["lesson_notes"])
        
        prompt = f"""
        You are an expert educator. Generate educational content for:
        
        SUBJECT: {request.subject_name}
        GRADE LEVEL: {request.grade_level}
        TOPIC: {request.topic_title}
        
        EXISTING TOPIC CONTENT:
        {json.dumps(request.topic_content, indent=2)}
        
        TASK:
        {instruction}
        
        Make the content appropriate for Grade {request.grade_level} students.
        Be comprehensive but age-appropriate.
        """
        
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{
                "role": "user",
                "content": prompt
            }],
            temperature=0.4,
            response_format={"type": "json_object"}
        )
        
        generated_content = json.loads(response.choices[0].message.content)
        
        return {
            "topic_id": request.topic_id,
            "resource_type": request.resource_type,
            "content": generated_content,
            "generated_at": "now"
        }
        
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Error generating resource: {str(e)}"
        )


# ==========================================
# MCQ GENERATION
# ==========================================

class MCQGenerationRequest(BaseModel):
    topic_id: str
    topic_title: str
    topic_content: dict
    question_count: int = 10
    difficulty_mix: dict = {"easy": 3, "medium": 5, "hard": 2}
    mcq_type: str = "daily"  # daily, weekly, chapter

@app.post("/api/mcq/generate")
async def generate_mcqs(request: MCQGenerationRequest):
    """
    Generate unique MCQ sets for a topic using the hybrid approach.
    Each generation produces different questions to avoid repetition.
    """
    if not client:
        raise HTTPException(status_code=500, detail="OpenAI API not configured")
    
    try:
        prompt = f"""
        Generate {request.question_count} UNIQUE multiple choice questions for:
        
        TOPIC: {request.topic_title}
        TYPE: {request.mcq_type} assessment
        
        CONTENT TO COVER:
        {json.dumps(request.topic_content, indent=2)}
        
        DIFFICULTY DISTRIBUTION:
        - Easy: {request.difficulty_mix.get('easy', 3)} questions
        - Medium: {request.difficulty_mix.get('medium', 5)} questions
        - Hard: {request.difficulty_mix.get('hard', 2)} questions
        
        REQUIREMENTS:
        1. Each question must be unique and test different aspects
        2. All 4 options should be plausible
        3. Include detailed explanations
        4. Vary question styles (direct, application, reasoning)
        
        Format as JSON:
        {{
            "questions": [
                {{
                    "id": 1,
                    "question": "...",
                    "options": {{"A": "...", "B": "...", "C": "...", "D": "..."}},
                    "correct_answer": "B",
                    "explanation": "...",
                    "difficulty": "medium",
                    "skill_tested": "application"
                }}
            ],
            "metadata": {{
                "total_questions": {request.question_count},
                "difficulty_breakdown": {json.dumps(request.difficulty_mix)}
            }}
        }}
        """
        
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{
                "role": "user",
                "content": prompt
            }],
            temperature=0.7,  # Higher for variety
            response_format={"type": "json_object"}
        )
        
        return json.loads(response.choices[0].message.content)
        
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Error generating MCQs: {str(e)}"
        )


if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
